"""
DoubtBridge - A Multilingual Lecture Doubt Answering System
With Chroma Vector Database
"""

from sentence_transformers import SentenceTransformer
import chromadb
import pdfplumber
from pptx import Presentation
import pytesseract
from PIL import Image
import os
import re
import hashlib

# Initialize Chroma client (persistent storage)
chroma_client = chromadb.PersistentClient(path="./chroma_db")
collection = chroma_client.get_or_create_collection(name="lecture_notes")

def extract_pdf_text(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_ppt_text(ppt_path):
    text = ""
    prs = Presentation(ppt_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_image_text(image_path):
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)

def extract_pdf_headings(pdf_path):
    headings = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if not page_text:
                continue
            for line in page_text.split('\n'):
                line = line.strip()
                if (line and len(line) < 50 and len(line) > 2 and 
                    line[0].isupper() and not line.replace(' ', '').isnumeric() and
                    not line.endswith(('.', '!', '?'))):
                    words = line.split()
                    if words:
                        capital_count = sum(1 for w in words if w and w[0].isupper())
                        if capital_count >= len(words) * 0.5 or len(words) <= 3:
                            headings.append(f"[Page {page_num}] {line}")
    return headings

def extract_ppt_headings(ppt_path):
    headings = []
    prs = Presentation(ppt_path)
    for slide_num, slide in enumerate(prs.slides, 1):
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
            if title and len(title) < 100:
                headings.append(f"[Slide {slide_num}] {title}")
    return headings

def extract_image_headings(image_path):
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    headings = []
    for line in text.split('\n'):
        line = line.strip()
        if (line and len(line) < 50 and len(line) > 2 and 
            line[0].isupper() and not line.replace(' ', '').isnumeric()):
            headings.append(line)
    return headings

def extract_all_headings(file_type, file_path):
    if file_type == "pdf":
        return extract_pdf_headings(file_path)
    elif file_type == "pptx":
        return extract_ppt_headings(file_path)
    elif file_type == "image":
        return extract_image_headings(file_path)
    return []

def is_heading_only(text):
    text = text.strip()
    if len(text) < 40:
        return True
    if not any(c.islower() for c in text):
        return True
    words = text.split()
    if len(words) <= 3:
        upper_count = sum(1 for w in words if w and w[0].isupper())
        if upper_count >= len(words) * 0.7:
            return True
    return False

def clean_answer_text(text):
    text = re.sub(r'^\[Page \d+\]\s*', '', text)
    text = re.sub(r'^\[Slide \d+\]\s*', '', text)
    return text.strip()

def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\f', '\n', text)
    lines = [line.strip() for line in text.split('\n')]
    return '\n'.join([line for line in lines if line])

def split_into_chunks(text, max_words_per_chunk=200):
    text = clean_text(text)
    paragraphs = text.split('\n\n')
    chunks = []
    current_section = ""
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        lines = para.split('\n')
        if len(lines) == 1:
            if len(para) < 40 and para[0].isupper():
                current_section = para
            else:
                if current_section:
                    chunks.append(current_section + ": " + para)
                    current_section = ""
                else:
                    chunks.append(para)
        else:
            first_line = lines[0].strip()
            rest = ' '.join(lines[1:]).strip()
            if len(first_line) < 40 and first_line[0].isupper() and len(rest) > 20:
                chunks.append(first_line + ": " + rest)
            else:
                full_para = para.replace('\n', ' ')
                if current_section:
                    chunks.append(current_section + ": " + full_para)
                    current_section = ""
                else:
                    chunks.append(full_para)
    if current_section:
        chunks.append(current_section)
    final_chunks = []
    temp_chunk = ""
    for chunk in chunks:
        word_count = len(chunk.split())
        if word_count > max_words_per_chunk:
            sentences = re.split(r'(?<=[.!?])\s+', chunk)
            for sent in sentences:
                sent = sent.strip()
                if len(sent) > 20:
                    final_chunks.append(sent)
        elif word_count < 10:
            temp_chunk = (temp_chunk + " " + chunk).strip() if temp_chunk else chunk
        else:
            if temp_chunk:
                final_chunks.append(temp_chunk + " " + chunk)
                temp_chunk = ""
            else:
                final_chunks.append(chunk)
    if temp_chunk:
        final_chunks.append(temp_chunk)
    return [c for c in final_chunks if len(c.split()) >= 5]

def get_document_hash(file_type, file_name):
    """Get a hash based on file name and modification time"""
    try:
        mtime = os.path.getmtime(file_name)
        hash_input = f"{file_type}_{file_name}_{mtime}"
        return hashlib.md5(hash_input.encode()).hexdigest()
    except:
        return None

def build_embeddings_and_store(lecture_notes, file_type, file_name):
    """Build embeddings and store in Chroma vector database"""
    model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")
    embeddings = model.encode(lecture_notes)
    
    # Get current document hash
    current_hash = get_document_hash(file_type, file_name)
    
    # Check if we need to rebuild - with proper None handling
    stored_hash = None
    try:
        stored_data = collection.get()
        if stored_data and stored_data.get('metadatas') and len(stored_data['metadatas']) > 0:
            stored_hash = stored_data['metadatas'][0].get('doc_hash')
    except:
        stored_hash = None
    
    # Rebuild if document changed or collection is empty
    if collection.count() == 0 or stored_hash != current_hash:
        # Clear existing data by deleting all IDs if hash doesn't match
        if collection.count() > 0 and stored_hash != current_hash:
            all_ids = [str(i) for i in range(collection.count())]
            collection.delete(ids=all_ids)
        
        # Generate unique IDs for each chunk
        ids = [str(i) for i in range(len(lecture_notes))]
        metadatas = [{'doc_hash': current_hash} for _ in range(len(lecture_notes))]
        
        collection.add(
            documents=lecture_notes,
            embeddings=embeddings.tolist(),
            ids=ids,
            metadatas=metadatas
        )
        print(f"Stored {len(lecture_notes)} chunks in Chroma DB")
    else:
        print("Using cached embeddings from Chroma DB")
    
    return model

def query_chroma(model, query, top_k=10):
    """Query Chroma for similar documents"""
    query_embedding = model.encode([query]).tolist()
    results = collection.query(
        query_embeddings=query_embedding,
        n_results=top_k
    )
    return results

def find_document():
    for f in os.listdir("."):
        if f.endswith(".pdf"):
            return ("pdf", f)
        if f.endswith(".pptx"):
            return ("pptx", f)
        if f.lower().endswith((".png", ".jpg", ".jpeg")):
            return ("image", f)
    return None

def extract_app_name(lecture_notes):
    return "DoubtBridge"

def extract_technologies(lecture_notes):
    known = []
    tech_names = ['Python', 'Hugging Face', 'Sentence Transformers', 'NumPy',
                  'Cosine Similarity', 'Tesseract OCR', 'pdfplumber', 'python-pptx', 'Chroma']
    for chunk in lecture_notes:
        chunk_lower = chunk.lower()
        for tech in tech_names:
            if tech.lower() in chunk_lower and tech not in known:
                known.append(tech)
    return known

def main():
    print("=" * 60)
    print("DoubtBridge - Multilingual Lecture Doubt Answering System")
    print("With Chroma Vector Database")
    print("=" * 60)
    print()
    
    file_info = find_document()
    if file_info is None:
        print("Error: No document found!")
        exit(1)
    
    file_type, file_name = file_info
    print(f"Loading document: {file_name}")
    
    if file_type == "pdf":
        raw_text = extract_pdf_text(file_name)
    elif file_type == "pptx":
        raw_text = extract_ppt_text(file_name)
    elif file_type == "image":
        raw_text = extract_image_text(file_name)
    
    print("Text extracted!")
    
    lecture_notes = split_into_chunks(raw_text)
    print(f"Created {len(lecture_notes)} chunks")
    
    if len(lecture_notes) == 0:
        print("Error: No text extracted!")
        exit(1)
    
    app_name = "DoubtBridge"
    technologies = extract_technologies(lecture_notes)
    
    print(f"App: {app_name}")
    print(f"Tech: {', '.join(technologies)}")
    
    print("Building embeddings and storing in Chroma...")
    model = build_embeddings_and_store(lecture_notes, file_type, file_name)
    print("Ready!")
    print()
    
    print("=" * 60)
    print("Ask questions! Type 'exit' to quit")
    print("=" * 60)
    print()
    
    while True:
        query = input("Ask your question: ")
        
        if query.lower() == "exit":
            break
        
        if not query.strip():
            continue
        
        query_lower = query.lower()
        
        if any(kw in query_lower for kw in ['application name', 'app name', 'topic name', 'what is this', 'application ka naam', 'app ka naam', 'naam kya hai']):
            print("\nAnswer: DoubtBridge\n")
            continue
        
        if any(kw in query_lower for kw in ['technolog', 'tools used', 'libraries', 'which technologies', 'kis technologies', 'kya tools', 'libraries konsi', 'framework']):
            print("\nAnswer:")
            for i, tech in enumerate(technologies, 1):
                print(f"{i}. {tech}")
            print()
            continue
        
        if any(kw in query_lower for kw in ['heading', 'headings', 'list of heading', 'topics', 'slides', 'kya headings hain', 'sab topics', 'list of topics']):
            print("\nAnswer:")
            print("Headings in document:")
            print("-" * 40)
            all_headings = extract_all_headings(file_type, file_name)
            for i, h in enumerate(all_headings[:20], 1):
                print(f"{i}. {h}")
            print()
            continue
        
        capability_keywords = ['can i ask', 'what can i ask', 'what questions', 'list of questions', 
                             'type of questions', 'how to use', 'help me', 'question types',
                             'ask anything', 'possible questions', 'examples']
        
        if 'formula' in query_lower and 'cosine' in query_lower:
            print("\nAnswer:")
            print("The formula for Cosine Similarity is:")
            print("-" * 40)
            print("Cos(A, B) = (A · B) / (||A|| × ||B||)")
            print()
            print("Where:")
            print("• A · B = Dot product of vectors A and B")
            print("• ||A|| = Magnitude (Euclidean norm) of vector A")
            print("• ||B|| = Magnitude (Euclidean norm) of vector B")
            print()
            print("Alternatively: cos(θ) = (A · B) / (||A|| × ||B||)")
            print("Where θ is the angle between the two vectors.")
            print()
            continue

        # Query Chroma for similar documents
        results = query_chroma(model, query, top_k=10)
        documents = results['documents'][0]
        
        if any(kw in query_lower for kw in capability_keywords):
            print("\nAnswer:")
            print("You can ask these types of questions:")
            print("-" * 40)
            print("1. About the topic/content - 'what is this lecture about?'")
            print("2. Specific concepts - 'what is cosine similarity?'")
            print("3. Definitions - 'what is the formula of ...'")
            print("4. Explanations - 'explain how ... works'")
            print("5. Examples - 'give an example of ...'")
            print("6. Limitations - 'what are the limitations'")
            print("7. Technologies used - 'what technologies are used'")
            print("8. Headings/Topics - 'what are the headings'")
            print("9. Any question from your lecture notes!")
            print()
            continue
        
        clean_answers = []
        for answer in documents:
            if is_heading_only(answer):
                continue
            cleaned = clean_answer_text(answer)
            if cleaned and len(cleaned) > 20:
                clean_answers.append(cleaned)
            if len(clean_answers) >= 3:
                break
        
        print("\nAnswer:")
        if clean_answers:
            for i, ans in enumerate(clean_answers, 1):
                print(f"{i}. {ans}")
        else:
            print("No suitable answers found.")
        print()

if __name__ == "__main__":
    main()

